from __future__ import annotations

import io
import csv
import base64
import uuid
from collections import defaultdict
from html import escape
from importlib.resources import files
from pathlib import Path
from typing import Any, Iterable

from PIL import Image as PILImage, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import HRFlowable, Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from .releases import release_snapshot_sha256, validate_release_definition, validate_release_items

THEME_PALETTES = {
    "forest": {"primary": "#086650", "ink": "#17231F", "paper": "#F4F1E8", "card": "#FFFFFF"},
    "industrial": {"primary": "#C34A21", "ink": "#22272B", "paper": "#ECEBE7", "card": "#FFFFFF"},
    "midnight": {"primary": "#2E63C7", "ink": "#111827", "paper": "#E9EEF7", "card": "#FFFFFF"},
    "classic": {"primary": "#8A6A2F", "ink": "#211D17", "paper": "#F5F0E5", "card": "#FFFDF8"},
}
CATALOG_THEMES = tuple(THEME_PALETTES)
NATSUKI_TITLE_FONT = "BarlowCondensed-Bold"
NATSUKI_BODY_FONT = "DMSans-Regular"
NATSUKI_BODY_BOLD_FONT = "DMSans-Bold"
MINIMUM_CATALOG_FONT_SIZE = 12


def _register_natsuki_fonts() -> None:
    font_root = files("perfect_catalog").joinpath("assets/brands/natsuki/fonts")
    for name, filename in (
        (NATSUKI_TITLE_FONT, "BarlowCondensed-Bold.ttf"),
        (NATSUKI_BODY_FONT, "DMSans-Regular.ttf"),
        (NATSUKI_BODY_BOLD_FONT, "DMSans-Bold.ttf"),
    ):
        if name not in pdfmetrics.getRegisteredFontNames():
            pdfmetrics.registerFont(TTFont(name, str(font_root.joinpath(filename))))


def _theme(config: dict[str, Any]) -> dict[str, str]:
    profile = config.get("visual_profile") or {}
    if profile:
        return {"primary": str(profile.get("primary_color") or "#E30613"), "secondary": str(profile.get("secondary_color") or "#12355B"), "ink": str(profile.get("ink_color") or "#111111"), "paper": str(profile.get("paper_color") or "#FFFFFF"), "card": "#FFFFFF"}
    return THEME_PALETTES.get(str(config.get("theme") or "forest"), THEME_PALETTES["forest"])


def _visual(config: dict[str, Any]) -> dict[str, Any]:
    return dict(config.get("visual_profile") or {})


def _logo_path(
    config: dict[str, Any], bundle_dir: Path | None = None, *,
    company: bool = False, raster_only: bool = False,
) -> Path | None:
    profile = _visual(config)
    filename = (profile.get("company") or {}).get("packaged_logo_path") if company else profile.get("packaged_logo_path")
    if filename and bundle_dir:
        candidate = (bundle_dir.resolve() / str(filename)).resolve()
        if candidate.is_relative_to(bundle_dir.resolve()) and candidate.is_file():
            if not raster_only or candidate.suffix.lower() in {".png", ".jpg", ".jpeg"}:
                return candidate
    if company:
        return None
    if str(_visual(config).get("logo_asset_key") or "") != "brands/natsuki/logo.svg": return None
    path = files("perfect_catalog").joinpath("assets/brands/natsuki/logo.png")
    return Path(str(path)) if path.is_file() else None


def _layout(config: dict[str, Any]) -> tuple[int, int]:
    layouts = {"T4": (2, 4), "T2": (1, 2), "T1": (1, 1), "TABLE": (1, 10)}
    profile = str(config.get("template_profile") or "").upper()
    if profile in layouts: return layouts[profile]
    columns = max(1, min(3, int(config.get("columns_per_row", 2))))
    return columns, {1: 2, 2: 4, 3: 6}[columns]


def export_rows_from_release(release: dict[str, Any], items: Iterable[dict[str, Any]]) -> list[dict[str, Any]]:
    """Adapta exclusivamente un release completo cuya integridad criptográfica fue revalidada."""
    materialized = sorted(list(items), key=lambda item: item["item_order"])
    validate_release_definition(release["definition"], len(materialized))
    validate_release_items(materialized)
    calculated = release_snapshot_sha256(uuid.UUID(str(release["brand_id"])), str(release["version"]), release["definition"], materialized)
    if calculated != release["snapshot_sha256"]:
        raise ValueError("Los items no coinciden con snapshot_sha256 del release.")
    rows = [dict(item["snapshot_data"]) for item in materialized]
    for row in rows:
        for excluded in ("quantity_available", "quantity_on_hand", "uom_original", "currency", "price", "price_two"):
            row.pop(excluded, None)
    return rows


def _groups(
    rows: list[dict[str, Any]], key: str, secondary_key: str | None = None
) -> list[tuple[str, list[dict[str, Any]]]]:
    grouped: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for row in rows:
        primary_values = (
            row.get("vehicle_makes") or ["Sin marca vehicular"]
            if key == "vehicle_make" else [row.get(key) or "Sin categoría"]
        )
        secondary_values = (
            row.get("vehicle_makes") or ["Sin marca vehicular"]
            if secondary_key == "vehicle_make" else [row.get(secondary_key) or "Sin subgrupo"]
            if secondary_key else [""]
        )
        for primary in primary_values:
            for secondary in secondary_values:
                label = f"{primary} · {secondary}" if secondary else str(primary)
                grouped[label].append(row)
    return list(grouped.items())


def _detail(row: dict[str, Any]) -> str:
    parts = [escape(str(row.get("name_original") or "")), f"Ref. {escape(str(row.get('internal_reference_original') or ''))}"]
    if row.get("piece_type") or row.get("category_path"):
        parts.append("Tipo: " + escape(str(row.get("piece_type") or row["category_path"])))
    if row.get("brand"):
        parts.append("Marca: " + escape(str(row["brand"])))
    if row.get("oem_references"):
        parts.append("OEM: " + escape(", ".join(map(str, row["oem_references"]))))
    if row.get("applications"):
        parts.append("Aplicaciones: " + escape("; ".join(str(value) for value in row["applications"])))
    if row.get("engine_types"):
        parts.append("Motor: " + escape(", ".join(map(str, row["engine_types"]))))
    return "<br/>".join(parts)


def _safe_bundle_image(row: dict[str, Any], bundle_dir: Path | None) -> Path | None:
    if bundle_dir is None or not row.get("image_path"):
        return None
    root = bundle_dir.resolve()
    candidate = (root / str(row["image_path"])).resolve()
    if not candidate.is_relative_to(root) or not candidate.is_file():
        return None
    return candidate


def _contained_size(
    source_width: float, source_height: float, max_width: float, max_height: float,
) -> tuple[float, float]:
    """Reduce proporcionalmente dentro de una caja; nunca recorta ni deforma."""
    if min(source_width, source_height, max_width, max_height) <= 0:
        raise ValueError("Las dimensiones de imagen deben ser positivas.")
    scale = min(max_width / source_width, max_height / source_height, 1.0)
    return source_width * scale, source_height * scale


def _optimized_raster(
    image_path: Path, max_width_px: int, max_height_px: int, *, quality: int = 84,
) -> io.BytesIO:
    """Creates a bounded display copy while leaving the approved original untouched."""
    with PILImage.open(image_path) as source:
        image = ImageOps.exif_transpose(source)
        image.thumbnail((max_width_px, max_height_px), PILImage.Resampling.LANCZOS)
        if image.mode not in {"RGB", "L"}:
            background = PILImage.new("RGB", image.size, "white")
            if "A" in image.getbands():
                background.paste(image, mask=image.getchannel("A"))
            else:
                background.paste(image.convert("RGB"))
            image = background
        elif image.mode == "L":
            image = image.convert("RGB")
        output = io.BytesIO()
        image.save(output, format="JPEG", quality=quality, optimize=True, progressive=True)
    output.seek(0)
    return output


def _pdf_cell(
    row: dict[str, Any], styles: Any, bundle_dir: Path | None,
    palette: dict[str, str], columns: int,
) -> list[Any]:
    contents: list[Any] = []
    image_path = _safe_bundle_image(row, bundle_dir)
    if image_path:
        try:
            available_width = (A4[0] - 3.2 * cm) / columns
            available_height = (4.4 if columns == 1 else 3.2) * cm
            optimized = _optimized_raster(
                image_path,
                max(1, round(available_width / 72 * 200)),
                max(1, round(available_height / 72 * 200)),
            )
            width, height = ImageReader(optimized).getSize()
            scale = min(available_width / width, (4.4 if columns == 1 else 3.2) * cm / height)
            optimized.seek(0)
            product_image = Image(optimized, width=width * scale, height=height * scale)
            product_image.hAlign = "CENTER"
            contents.extend([product_image, Spacer(1, .22 * cm)])
        except Exception as exc:
            reference = str(row.get("internal_reference_original") or "sin referencia")
            raise RuntimeError(f"No se pudo preparar la imagen PDF de {reference}.") from exc
    reference = escape(str(row.get("internal_reference_original") or "Sin referencia"))
    contents.append(Paragraph(reference, styles["CatalogReference"]))
    contents.append(Paragraph(escape(str(row.get("name_original") or "Sin nombre")), styles["CatalogProductTitle"]))
    detail_parts: list[str] = []
    if row.get("piece_type") or row.get("category_path"):
        detail_parts.append("<b>Tipo</b> · " + escape(str(row.get("piece_type") or row["category_path"])))
    if row.get("applications"):
        detail_parts.append("<b>Aplicaciones</b><br/>" + escape("; ".join(map(str, row["applications"]))))
    if row.get("engine_types"):
        detail_parts.append("<b>Motor</b> · " + escape(", ".join(map(str, row["engine_types"]))))
    if row.get("oem_references"):
        detail_parts.append("<b>OEM</b> · " + escape(", ".join(map(str, row["oem_references"]))))
    if detail_parts:
        contents.append(Paragraph("<br/>".join(detail_parts), styles["CatalogMeta"]))
    return contents


def generate_catalog_pdf(
    rows: list[dict[str, Any]], config: dict[str, Any] | None = None,
    *, bundle_dir: Path | None = None, release: dict[str, Any] | None = None,
) -> bytes:
    config = config or {}
    release = release or {}
    columns, page_capacity = _layout(config)
    title = str(config.get("title") or "Catálogo de productos")
    palette = _theme(config)
    _register_natsuki_fonts()
    styles = getSampleStyleSheet()
    cover_title_style = ParagraphStyle(
        "PerfectCatalogCoverTitle", parent=styles["Title"],
        textColor=colors.HexColor(palette["ink"]), fontName=NATSUKI_TITLE_FONT,
        fontSize=38, leading=41, alignment=0, spaceAfter=14,
    )
    cover_subtitle_style = ParagraphStyle(
        "PerfectCatalogCoverSubtitle", parent=styles["Heading2"],
        textColor=colors.HexColor(palette["ink"]), fontName=NATSUKI_BODY_FONT,
        fontSize=15, leading=20, alignment=0,
    )
    section_style = ParagraphStyle(
        "PerfectCatalogSection", parent=styles["Heading1"],
        textColor=colors.HexColor(palette["ink"]), fontName=NATSUKI_TITLE_FONT,
        fontSize=23, leading=27, spaceAfter=4,
    )
    styles.add(ParagraphStyle(
        "CatalogEyebrow", parent=styles["Normal"], textColor=colors.HexColor(palette["primary"]),
        fontName=NATSUKI_BODY_BOLD_FONT, fontSize=12, leading=21.6, spaceAfter=10,
    ))
    styles.add(ParagraphStyle(
        "CatalogReference", parent=styles["Normal"], textColor=colors.HexColor(palette["primary"]),
        fontName=NATSUKI_BODY_BOLD_FONT, fontSize=12, leading=21.6, spaceAfter=5,
    ))
    styles.add(ParagraphStyle(
        "CatalogProductTitle", parent=styles["Heading3"], textColor=colors.HexColor(palette["ink"]),
        fontName=NATSUKI_TITLE_FONT, fontSize=14,
        leading=22, spaceAfter=8,
    ))
    styles.add(ParagraphStyle(
        "CatalogMeta", parent=styles["BodyText"], textColor=colors.HexColor("#56645e"),
        fontName=NATSUKI_BODY_FONT, fontSize=12,
        leading=21.6,
    ))
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=A4, title=title, author="Perfect Trading",
        subject="Catálogo verificable de productos",
        leftMargin=1.35*cm, rightMargin=1.35*cm, topMargin=1.55*cm, bottomMargin=1.35*cm,
    )
    version = str(release.get("version") or "Edición de trabajo")
    checksum = str(release.get("snapshot_sha256") or "")
    cover_meta = Table(
        [[Paragraph("EDICIÓN", styles["CatalogEyebrow"]), Paragraph("PRODUCTOS", styles["CatalogEyebrow"]), Paragraph("IDENTIDAD", styles["CatalogEyebrow"])],
         [Paragraph(escape(version), styles["BodyText"]), Paragraph(str(len(rows)), styles["BodyText"]), Paragraph(escape(checksum[:16] or "Sin publicar"), styles["BodyText"]) ]],
        colWidths=[6.2 * cm, 3.2 * cm, 7.4 * cm],
    )
    cover_meta.setStyle(TableStyle([
        ("LINEABOVE", (0, 0), (-1, 0), 1, colors.HexColor(palette["primary"])),
        ("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 0),
        ("RIGHTPADDING", (0, 0), (-1, -1), 8), ("TOPPADDING", (0, 0), (-1, -1), 8),
    ]))
    story: list[Any] = [
        Spacer(1, 4.4 * cm), Paragraph("PERFECT TRADING · CATÁLOGO", styles["CatalogEyebrow"]),
        Paragraph(escape(title), cover_title_style),
        Paragraph(escape(str(config.get("subtitle") or "Selección técnica de productos")), cover_subtitle_style),
        Spacer(1, 2.2 * cm), cover_meta, PageBreak(),
    ]
    for section_index, (section, section_rows) in enumerate(_groups(
        rows, str(config.get("group_by") or "category_path"),
        str(config["group_by_secondary"]) if config.get("group_by_secondary") else None,
    )):
        cells = [_pdf_cell(row, styles, bundle_dir, palette, columns) for row in section_rows]
        chunks = [cells[index:index + page_capacity] for index in range(0, len(cells), page_capacity)]
        for chunk_index, chunk in enumerate(chunks):
            if section_index or chunk_index:
                story.append(PageBreak())
            continuation = " · CONTINUACIÓN" if chunk_index else ""
            story.extend([
                Paragraph(f"SECCIÓN {section_index + 1:02d}{continuation}", styles["CatalogEyebrow"]),
                Paragraph(escape(section), section_style),
                Paragraph(f"{len(section_rows)} productos", styles["CatalogMeta"]),
                HRFlowable(width="100%", thickness=1, color=colors.HexColor(palette["primary"]), spaceBefore=7, spaceAfter=12),
            ])
            grid = [chunk[index:index + columns] for index in range(0, len(chunk), columns)]
            if grid and len(grid[-1]) < columns:
                grid[-1].extend([""] * (columns - len(grid[-1])))
            table = Table(grid, colWidths=[(A4[0]-2.7*cm)/columns]*columns, repeatRows=0, hAlign="LEFT")
            table.setStyle(TableStyle([
                ("BACKGROUND", (0,0), (-1,-1), colors.HexColor(palette["card"])),
                ("BOX", (0,0), (-1,-1), .45, colors.HexColor("#d7ddd9")),
                ("INNERGRID", (0,0), (-1,-1), .35, colors.HexColor("#e3e7e4")),
                ("LINEABOVE", (0,0), (-1,0), 2, colors.HexColor(palette["primary"])),
                ("VALIGN", (0,0), (-1,-1), "TOP"), ("PADDING", (0,0), (-1,-1), 11),
            ]))
            story.extend([table, Spacer(1, .5*cm)])

    def decorate_cover(canvas: Any, document: Any) -> None:
        canvas.saveState()
        canvas.setTitle(title)
        canvas.setAuthor("Perfect Trading International")
        canvas.setCreator("Perfect Catalog")
        canvas.setSubject("Catálogo verificable de productos")
        canvas.setKeywords("catálogo, autopartes, Perfect Trading, productos")
        canvas.bookmarkPage("portada")
        canvas.addOutlineEntry(title, "portada", level=0, closed=False)
        canvas.setFillColor(colors.HexColor(palette["paper"]))
        canvas.rect(0, 0, A4[0], A4[1], stroke=0, fill=1)
        canvas.setFillColor(colors.HexColor(palette["primary"]))
        canvas.rect(0, A4[1] - 1.1 * cm, A4[0], 1.1 * cm, stroke=0, fill=1)
        canvas.circle(A4[0] - 2.4 * cm, A4[1] - 3.5 * cm, 1.25 * cm, stroke=0, fill=1)
        company_logo = _logo_path(config, bundle_dir, company=True, raster_only=True)
        brand_logo = _logo_path(config, bundle_dir, raster_only=True)
        profile = _visual(config)
        if brand_logo and profile.get("watermark_enabled", True):
            canvas.saveState()
            canvas.setFillAlpha(float(profile.get("watermark_opacity") or .05))
            canvas.drawImage(str(brand_logo), A4[0]-11*cm, 1.8*cm, width=9*cm, height=1.7*cm, preserveAspectRatio=True, mask="auto")
            canvas.restoreState()
        corner_logo = company_logo or brand_logo
        if corner_logo and profile.get("corner_logo_enabled", True):
            canvas.drawImage(str(corner_logo), 1.35*cm, A4[1]-2.3*cm, width=4.3*cm, height=.82*cm, preserveAspectRatio=True, mask="auto")
        canvas.restoreState()

    def decorate(canvas: Any, document: Any) -> None:
        canvas.saveState()
        canvas.setStrokeColor(colors.HexColor(palette["primary"]))
        canvas.setLineWidth(2)
        canvas.line(1.2 * cm, A4[1] - .8 * cm, A4[0] - 1.2 * cm, A4[1] - .8 * cm)
        canvas.setFillColor(colors.HexColor(palette["ink"]))
        canvas.setFont(NATSUKI_BODY_FONT, MINIMUM_CATALOG_FONT_SIZE)
        canvas.drawString(1.35 * cm, A4[1] - 1.15 * cm, title[:70])
        proof = " · ".join(value for value in (version, checksum[:16]) if value)
        canvas.drawString(1.2 * cm, .65 * cm, proof)
        canvas.drawRightString(A4[0] - 1.2 * cm, .65 * cm, f"Página {document.page}")
        logo = _logo_path(config, bundle_dir, raster_only=True)
        if logo and _visual(config).get("corner_logo_enabled", True):
            canvas.drawImage(str(logo), A4[0]-5.4*cm, A4[1]-1.45*cm, width=4.1*cm, height=.78*cm, preserveAspectRatio=True, mask="auto")
        canvas.restoreState()

    doc.build(story, onFirstPage=decorate_cover, onLaterPages=decorate)
    return buffer.getvalue()


def generate_catalog_pptx(
    rows: list[dict[str, Any]], config: dict[str, Any] | None = None,
    *, bundle_dir: Path | None = None, release: dict[str, Any] | None = None,
) -> bytes:
    config = config or {}
    release = release or {}
    columns, per_slide = _layout(config)
    title = str(config.get("title") or "Catálogo de productos")
    palette = _theme(config)
    primary_rgb = RGBColor.from_string(palette["primary"].lstrip("#"))
    prs = Presentation()
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.background.fill.solid(); cover.background.fill.fore_color.rgb = RGBColor.from_string(palette["paper"].lstrip("#"))
    cover.shapes.title.text = title
    cover.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_rgb
    proof = " · ".join(value for value in (
        str(config.get("subtitle") or ""), str(release.get("version") or ""),
        str(release.get("snapshot_sha256") or "")[:16],
    ) if value)
    cover.placeholders[1].text = proof
    cover.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(palette["ink"].lstrip("#"))
    logo = (_logo_path(config, bundle_dir, company=True, raster_only=True)
            or _logo_path(config, bundle_dir, raster_only=True))
    if logo and _visual(config).get("corner_logo_enabled", True):
        cover.shapes.add_picture(str(logo), Inches(9.2), Inches(.35), width=Inches(3.5))
    for section, section_rows in _groups(
        rows, str(config.get("group_by") or "category_path"),
        str(config["group_by_secondary"]) if config.get("group_by_secondary") else None,
    ):
        for start in range(0, len(section_rows), per_slide):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor.from_string(palette["paper"].lstrip("#"))
            heading = slide.shapes.add_textbox(Inches(.4), Inches(.2), Inches(12.5), Inches(.5))
            heading.text_frame.paragraphs[0].text = section
            heading.text_frame.paragraphs[0].font.size = Pt(22)
            heading.text_frame.paragraphs[0].font.bold = True
            heading.text_frame.paragraphs[0].font.color.rgb = primary_rgb
            for index, row in enumerate(section_rows[start:start+per_slide]):
                col, line = index % columns, index // columns
                width = 12.4 / columns
                box = slide.shapes.add_textbox(Inches(.4+col*width), Inches(.9+line*2.05), Inches(width-.15), Inches(1.85))
                box.fill.solid(); box.fill.fore_color.rgb = RGBColor.from_string(palette["card"].lstrip("#"))
                box.line.color.rgb = primary_rgb
                frame = box.text_frame; frame.clear()
                image_path = _safe_bundle_image(row, bundle_dir)
                if image_path:
                    try:
                        optimized = _optimized_raster(image_path, 300, 240)
                        source_width, source_height = ImageReader(optimized).getSize()
                        image_width, image_height = _contained_size(
                            source_width, source_height, 1.0, .8,
                        )
                        optimized.seek(0)
                        slide.shapes.add_picture(
                            optimized,
                            Inches(.55 + col * width + (1.0 - image_width) / 2),
                            Inches(1.05 + line * 2.05 + (.8 - image_height) / 2),
                            width=Inches(image_width), height=Inches(image_height),
                        )
                        frame.margin_left = Inches(1.15)
                    except Exception as exc:
                        reference = str(row.get("internal_reference_original") or "sin referencia")
                        raise RuntimeError(
                            f"No se pudo preparar la imagen PowerPoint de {reference}."
                        ) from exc
                p = frame.paragraphs[0]; p.text = str(row.get("internal_reference_original") or ""); p.font.bold = True; p.font.size = Pt(12)
                p = frame.add_paragraph(); p.text = str(row.get("name_original") or ""); p.font.size = Pt(12)
                if row.get("piece_type") or row.get("category_path") or row.get("brand"):
                    metadata = " · ".join(str(value) for value in (row.get("piece_type") or row.get("category_path"), row.get("brand")) if value)
                    p = frame.add_paragraph(); p.text = metadata; p.font.size = Pt(12)
                if row.get("oem_references"):
                    p = frame.add_paragraph(); p.text = "OEM: " + ", ".join(map(str, row["oem_references"])); p.font.size = Pt(12)
                if row.get("applications"):
                    p = frame.add_paragraph(); p.text = "Aplicaciones: " + "; ".join(map(str, row["applications"])); p.font.size = Pt(12)
                if row.get("engine_types"):
                    p = frame.add_paragraph(); p.text = "Motor: " + ", ".join(map(str, row["engine_types"])); p.font.size = Pt(12)
                for paragraph in frame.paragraphs:
                    paragraph.font.name = "DM Sans"
            if logo and _visual(config).get("corner_logo_enabled", True):
                slide.shapes.add_picture(str(logo), Inches(10.7), Inches(.15), width=Inches(2.0))
    output = io.BytesIO(); prs.save(output)
    return output.getvalue()


def generate_catalog_html(
    rows: list[dict[str, Any]], config: dict[str, Any] | None = None,
    *, release: dict[str, Any] | None = None, bundle_dir: Path | None = None,
    embed_images: bool = False,
) -> bytes:
    """Genera una edición digital portable; no consulta estado mutable ni ejecuta JavaScript."""
    config = config or {}
    release = release or {}
    columns, _ = _layout(config)
    palette = _theme(config)
    title = escape(str(config.get("title") or "Catálogo de productos"))
    subtitle = escape(str(config.get("subtitle") or ""))
    sections: list[str] = []
    lightboxes: list[str] = []
    grouped_rows = _groups(
        rows, str(config.get("group_by") or "category_path"),
        str(config["group_by_secondary"]) if config.get("group_by_secondary") else None,
    )
    navigation: list[str] = []
    for section_index, (section, section_rows) in enumerate(grouped_rows, 1):
        section_id = f"seccion-{section_index:02d}"
        navigation.append(f'<a href="#{section_id}"><span>{section_index:02d}</span>{escape(section)}</a>')
        cards: list[str] = []
        for card_index, row in enumerate(section_rows, 1):
            image = ""
            if row.get("image_path"):
                source = str(row["image_path"])
                if embed_images:
                    image_path = _safe_bundle_image(row, bundle_dir)
                    if image_path is None:
                        raise FileNotFoundError(
                            f"No se puede incrustar la imagen segura {source!r}."
                        )
                    optimized = _optimized_raster(
                        image_path, 1200, 900, quality=82,
                    )
                    source = (
                        "data:image/jpeg;base64,"
                        + base64.b64encode(optimized.read()).decode("ascii")
                    )
                image_alt = escape(
                    str(row.get("internal_reference_original") or row.get("name_original") or "Producto"),
                    quote=True,
                )
                image_source = escape(source, quote=True)
                lightbox_id = f"foto-{section_index:02d}-{card_index:03d}"
                image = (
                    f'<a class="photo" href="#{lightbox_id}" aria-label="Ampliar imagen de {image_alt}">'
                    f'<img src="{image_source}" alt="{image_alt}" loading="lazy" decoding="async">'
                    '<span class="zoom-hint" aria-hidden="true">Ampliar</span></a>'
                )
                lightboxes.append(
                    f'<figure class="lightbox" id="{lightbox_id}" tabindex="-1">'
                    f'<a class="lightbox-backdrop" href="#{section_id}" aria-label="Cerrar imagen ampliada"></a>'
                    '<div class="lightbox-dialog">'
                    f'<img src="{image_source}" alt="{image_alt}">'
                    f'<figcaption>{image_alt}</figcaption>'
                    f'<a class="lightbox-close" href="#{section_id}" aria-label="Cerrar imagen ampliada">Cerrar</a>'
                    '</div></figure>'
                )
            applications = escape("; ".join(map(str, row.get("applications") or [])))
            engines = escape(", ".join(map(str, row.get("engine_types") or [])))
            category = escape(str(row.get("category_path") or ""))
            brand = escape(str(row.get("brand") or ""))
            oem = escape(", ".join(map(str, row.get("oem_references") or [])))
            specifications = (
                (f'<div><dt>OEM</dt><dd>{oem}</dd></div>' if oem else "")
                + (f'<div><dt>Aplicaciones</dt><dd>{applications}</dd></div>' if applications else "")
                + (f'<div><dt>Motor</dt><dd>{engines}</dd></div>' if engines else "")
            )
            cards.append(
                '<article class="product">' + image
                + f'<code>{escape(str(row.get("internal_reference_original") or "Sin referencia"))}</code>'
                + f'<h3>{escape(str(row.get("name_original") or "Sin nombre"))}</h3>'
                + (f'<p class="meta">{category}{" · " if category and brand else ""}{brand}</p>' if category or brand else "")
                + (f'<dl class="specifications">{specifications}</dl>' if specifications else "") + "</article>"
            )
        sections.append(
            f'<section id="{section_id}"><header><h2>{escape(section)}</h2><span>{len(section_rows)} productos</span></header>'
            f'<div class="products">{"".join(cards)}</div></section>'
        )
    checksum = escape(str(release.get("snapshot_sha256") or ""))
    version = escape(str(release.get("version") or ""))
    def logo_uri(path: Path | None) -> str:
        if not path: return ""
        media = {".svg": "image/svg+xml", ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg"}.get(path.suffix.lower(), "image/png")
        return f"data:{media};base64," + base64.b64encode(path.read_bytes()).decode("ascii")
    brand_logo_path = _logo_path(config, bundle_dir)
    company_logo_uri = logo_uri(_logo_path(config, bundle_dir, company=True) or brand_logo_path)
    brand_logo_uri = logo_uri(brand_logo_path)
    html = f"""<!doctype html>
<html lang="es"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<meta name="generator" content="Perfect Catalog"><meta name="release-sha256" content="{checksum}">
<title>{title}</title><style>
:root{{--ink:{palette['ink']};--forest:{palette['primary']};--paper:{palette['paper']};--card:{palette['card']};--line:#d9d5c9;--muted:#65716b}}*{{box-sizing:border-box}}html{{scroll-behavior:smooth}}body{{margin:0;color:var(--ink);background:var(--paper);font:15px/1.55 Arial,sans-serif}}main{{max-width:1280px;margin:auto;padding:clamp(24px,5vw,72px)}}.hero{{position:relative;min-height:48vh;display:grid;align-content:end;padding:8vw clamp(0px,2vw,28px) 4vw;border-bottom:4px solid var(--ink)}}.hero:before{{content:"";position:absolute;top:12%;right:2%;width:clamp(90px,14vw,190px);aspect-ratio:1;border:1px solid var(--forest);border-radius:50%;opacity:.22}}.hero small{{color:var(--forest);font-weight:800;letter-spacing:.16em;text-transform:uppercase}}h1{{position:relative;max-width:900px;margin:.2em 0;font:500 clamp(44px,8vw,104px)/.9 Georgia,serif;letter-spacing:-.035em}}.hero p{{max-width:700px;font-size:18px}}.contents{{display:flex;gap:8px;padding:20px 0;border-bottom:1px solid var(--line);overflow-x:auto;scrollbar-width:thin}}.contents a{{min-height:44px;display:inline-flex;gap:9px;align-items:center;flex:0 0 auto;padding:8px 13px;border:1px solid var(--line);border-radius:999px;color:var(--ink);background:var(--card);text-decoration:none}}.contents a:hover,.contents a:focus-visible{{border-color:var(--forest)}}.contents span{{color:var(--forest);font-weight:800}}section{{scroll-margin-top:18px;padding:clamp(38px,6vw,72px) 0}}section>header{{display:flex;justify-content:space-between;gap:20px;align-items:end;border-bottom:1px solid var(--line)}}h2{{margin:.25em 0;font:500 clamp(27px,4vw,48px) Georgia,serif;letter-spacing:-.02em}}section>header span{{padding-bottom:1.2em;color:var(--muted)}}.products{{display:grid;grid-template-columns:repeat({columns},minmax(0,1fr));gap:20px;padding-top:24px}}.product{{min-width:0;padding:20px;background:var(--card);border:1px solid var(--line);border-radius:14px;box-shadow:0 10px 30px rgba(20,42,34,.06);overflow:hidden}}.photo{{height:210px;margin:-20px -20px 20px;padding:10px;background:#f8f8f5;display:grid;place-items:center;overflow:hidden;border-bottom:1px solid var(--line)}}.photo img{{display:block;width:100%;height:100%;object-fit:contain;object-position:center center}}code{{color:var(--forest);font-weight:800;letter-spacing:.035em}}h3{{margin:.5em 0;font:500 22px/1.15 Georgia,serif}}.meta{{color:var(--muted);font-size:13px}}.specifications{{display:grid;gap:8px;margin:15px 0 0}}.specifications div{{display:grid;grid-template-columns:minmax(92px,.34fr) 1fr;gap:10px;padding-top:8px;border-top:1px solid var(--line)}}.specifications dt{{color:var(--forest);font-weight:800}}.specifications dd{{margin:0;overflow-wrap:anywhere}}.proof{{padding:28px 0;border-top:1px solid var(--line);overflow-wrap:anywhere;color:var(--muted);font-size:12px}}@media(max-width:760px){{html{{scroll-behavior:auto}}.products{{grid-template-columns:1fr}}.hero{{min-height:38vh}}section>header{{align-items:start;flex-direction:column;gap:0}}section>header span{{padding-bottom:1em}}}}@media(prefers-reduced-motion:reduce){{html{{scroll-behavior:auto}}}}@media print{{@page{{size:A4;margin:12mm}}body{{background:#fff}}main{{max-width:none;padding:0}}.hero{{min-height:245mm;break-after:page}}.contents{{display:none}}section{{break-before:page;padding:0}}.product{{break-inside:avoid;box-shadow:none}}.products{{gap:6mm}}}}
/* Visor ampliado sin JavaScript: mantiene el catálogo autónomo y portable. */
.photo{{position:relative;color:var(--ink);text-decoration:none;cursor:zoom-in}}.zoom-hint{{position:absolute;right:12px;bottom:12px;padding:6px 10px;border-radius:999px;color:#fff;background:rgba(17,30,25,.78);font-size:12px;font-weight:800;opacity:0;transform:translateY(4px);transition:.18s ease}}.photo:hover .zoom-hint,.photo:focus-visible .zoom-hint{{opacity:1;transform:none}}.lightbox{{position:fixed;inset:0;z-index:1000;display:none;margin:0;padding:clamp(18px,4vw,48px)}}.lightbox:target{{display:grid;place-items:center}}.lightbox-backdrop{{position:absolute;inset:0;background:rgba(8,15,12,.9)}}.lightbox-dialog{{position:relative;z-index:1;display:grid;grid-template-columns:1fr auto;grid-template-rows:minmax(0,1fr) auto;gap:12px;width:min(94vw,1400px);height:min(92vh,1000px);padding:16px;border-radius:16px;background:var(--card);box-shadow:0 24px 80px rgba(0,0,0,.45)}}.lightbox-dialog img{{grid-column:1/-1;width:100%;height:100%;min-height:0;object-fit:contain;object-position:center;background:#f8f8f5}}.lightbox-dialog figcaption{{align-self:center;overflow-wrap:anywhere}}.lightbox-close{{align-self:center;min-height:44px;padding:10px 16px;border:1px solid var(--line);border-radius:999px;color:var(--ink);font-weight:800;text-decoration:none}}@media(max-width:760px){{.zoom-hint{{opacity:1;transform:none}}}}@media(prefers-reduced-motion:reduce){{.zoom-hint{{transition:none}}}}@media print{{.lightbox{{display:none!important}}}}
</style></head><body><main><header class="hero"><small>Perfect Trading · edición {version}</small><h1>{title}</h1><p>{subtitle}</p></header><nav class="contents" aria-label="Secciones del catálogo">{''.join(navigation)}</nav>{''.join(sections)}<footer class="proof">Release SHA-256: {checksum}</footer></main>{''.join(lightboxes)}</body></html>"""
    brand_css = """@font-face{font-family:'DM Sans';src:url(data:font/ttf;base64,%s)}@font-face{font-family:'Barlow Condensed';src:url(data:font/ttf;base64,%s);font-weight:700}body{font-family:'DM Sans',sans-serif;font-size:16px;line-height:1.8}h1,h2,h3{font-family:'Barlow Condensed',sans-serif;font-weight:700}.meta,.proof{font-size:16px}.brand-logo{position:absolute;right:2rem;top:2rem;width:min(260px,35vw);z-index:2}.watermark{position:absolute;right:5%%;bottom:8%%;width:55%%;opacity:.05;pointer-events:none}""" % (
        base64.b64encode(files("perfect_catalog").joinpath("assets/brands/natsuki/fonts/DMSans-Regular.ttf").read_bytes()).decode("ascii"),
        base64.b64encode(files("perfect_catalog").joinpath("assets/brands/natsuki/fonts/BarlowCondensed-Bold.ttf").read_bytes()).decode("ascii"),
    )
    html = html.replace("</style>", brand_css + "</style>", 1)
    if company_logo_uri or brand_logo_uri:
        marks = ((f'<img class="brand-logo" src="{company_logo_uri}" alt="Perfect Trading">' if company_logo_uri else "")
                 + (f'<img class="watermark" src="{brand_logo_uri}" alt="">' if brand_logo_uri else ""))
        html = html.replace('<header class="hero">', '<header class="hero">' + marks, 1)
    return html.encode("utf-8")


def generate_indesign_datamerge_csv(rows: list[dict[str, Any]]) -> bytes:
    """CSV UTF-8 BOM para Data Merge, sin fórmulas interpretables al abrirlo en hojas de cálculo."""
    def cell(value: object) -> str:
        if isinstance(value, (list, tuple)):
            text = "; ".join(str(item) for item in value)
        else:
            text = "" if value is None else str(value)
        return "'" + text if text.startswith(("=", "+", "-", "@")) else text

    stream = io.StringIO(newline="")
    writer = csv.writer(stream, lineterminator="\r\n")
    writer.writerow(("reference", "name", "piece_type", "category", "brand", "vehicle_make", "applications", "engine_types", "oem_references", "@image"))
    for row in rows:
        writer.writerow((
            cell(row.get("internal_reference_original")), cell(row.get("name_original")),
            cell(row.get("piece_type")), cell(row.get("category_path")), cell(row.get("brand")),
            cell(row.get("vehicle_make") or row.get("vehicle_makes")),
            cell(row.get("applications")), cell(row.get("engine_types")), cell(row.get("oem_references")),
            cell(row.get("image_path")),
        ))
    return stream.getvalue().encode("utf-8-sig")
